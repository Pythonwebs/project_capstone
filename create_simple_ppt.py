"""
Simple PowerPoint Generator using python-pptx
Generates a basic 10-slide presentation for ServiceNow Incident Management System
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("Installing python-pptx...")
    import subprocess
    subprocess.check_call(["python", "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BLUE = RGBColor(31, 78, 121)
LIGHT_BLUE = RGBColor(68, 114, 196)
GREEN = RGBColor(76, 175, 80)
RED = RGBColor(244, 67, 54)
GRAY = RGBColor(89, 89, 89)
WHITE = RGBColor(255, 255, 255)

def add_title_slide(title, subtitle=""):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = LIGHT_BLUE
        p.alignment = PP_ALIGN.CENTER

def add_content_slide(title, bullets):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Header
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.color.rgb = DARK_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.7))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = GRAY
        p.space_before = Pt(6)
        p.space_after = Pt(6)

# SLIDE 1: Title
add_title_slide(
    "ServiceNow Incident Management",
    "Full-Stack Web Application"
)

# SLIDE 2: Overview
add_content_slide("Project Overview", [
    "📋 Web-based platform to manage ServiceNow incidents",
    "🔧 Complete CRUD operations (Create, Read, Update, Delete)",
    "🔐 OAuth 2.0 PKCE authentication",
    "🎨 Responsive UI with dark mode support",
    "⚡ Real-time synchronization with ServiceNow",
    "✅ Production-ready with error handling"
])

# SLIDE 3: Tech Stack
add_content_slide("Technology Stack", [
    "Frontend: React 19 + Vite + Material-UI",
    "Backend: Node.js + Express.js",
    "Authentication: OAuth 2.0 PKCE",
    "API Client: Axios",
    "Database: ServiceNow Instance",
    "Styling: Material-UI Components + Dark Mode"
])

# SLIDE 4: Architecture
add_content_slide("System Architecture", [
    "┌─────────────────────────────────────┐",
    "│  React Frontend (http://localhost:5173)",
    "└──────────────────┬──────────────────┘",
    "                   │",
    "┌──────────────────▼──────────────────┐",
    "│  Express Backend (http://localhost:3001)",
    "└──────────────────┬──────────────────┘",
    "                   │",
    "┌──────────────────▼──────────────────┐",
    "│  ServiceNow Instance (OAuth + REST API)",
    "└─────────────────────────────────────┘"
])

# SLIDE 5: Features
add_content_slide("Key Features", [
    "✓ Create Incidents - short_description, impact, urgency",
    "✓ View Incidents - grid layout with incident cards",
    "✓ Edit Incidents - update all incident details",
    "✓ Delete Incidents - with confirmation dialog",
    "✓ Dark Mode - toggle light/dark theme",
    "✓ Error Handling - user-friendly alerts",
    "✓ Session Management - secure OAuth flow"
])

# SLIDE 6: Frontend
add_content_slide("Frontend Components", [
    "Home.jsx - Main incident management interface",
    "  • Incident list display",
    "  • Create/Edit/Delete dialogs",
    "  • Success notifications",
    "",
    "App.jsx - Navigation and header",
    "  • AppBar with theme toggle",
    "  • Route management",
    "",
    "AuthProvider.jsx - Authentication context"
])

# SLIDE 7: Backend
add_content_slide("Backend Implementation", [
    "OAuth 2.0 PKCE Flow:",
    "  • Secure authentication with ServiceNow",
    "  • Token refresh on expiration",
    "",
    "API Endpoints:",
    "  • GET /api/incidents - List all",
    "  • POST /api/incidents - Create",
    "  • PUT /api/incidents/:id - Update",
    "  • DELETE /api/incidents/:id - Delete"
])

# SLIDE 8: Installation
add_content_slide("Setup Instructions", [
    "Backend Setup:",
    "  $ cd BFF",
    "  $ npm install",
    "  $ npm start  (port 3001)",
    "",
    "Frontend Setup:",
    "  $ cd client",
    "  $ npm install",
    "  $ npm run dev  (port 5173)",
    "",
    "Environment: Configure .env with OAuth credentials"
])

# SLIDE 9: API Example
add_content_slide("API Example: Create Incident", [
    "Request:",
    "  POST /api/incidents",
    "  {",
    '    "short_description": "System down",',
    '    "impact": 1,',
    '    "urgency": 1',
    "  }",
    "",
    "Response:",
    '  { "result": { "sys_id": "...", "number": "INC001" } }'
])

# SLIDE 10: Conclusion
add_title_slide(
    "Thank You!",
    "Questions? | GitHub: project_capstone"
)

# Save
output_path = r'c:\WEB TECH\project_capstone\ServiceNow_Incident_Management.pptx'
prs.save(output_path)
print(f"✅ PowerPoint presentation created successfully!")
print(f"📁 File: {output_path}")
print(f"📊 Slides: 10")
